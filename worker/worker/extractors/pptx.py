from __future__ import annotations
from io import BytesIO

from pptx import Presentation

from . import NormalizedDoc, PageContent


def extract(source_id: str, raw: bytes) -> NormalizedDoc:
    pres = Presentation(BytesIO(raw))
    pages: list[PageContent] = []
    for i, slide in enumerate(pres.slides):
        chunks = []
        title = None
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if not t:
                    continue
                if shape == slide.shapes.title and not title:
                    title = t
                chunks.append(t)
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        text = "\n".join(chunks) + (("\n[Speaker notes]\n" + notes) if notes else "")
        pages.append(
            PageContent(page_num=i + 1, text=text, section_title=title)
        )
    return NormalizedDoc(source_id=source_id, pages=pages)
